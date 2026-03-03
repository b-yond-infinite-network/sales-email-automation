from langgraph.checkpoint.memory import MemorySaver
from langgraph.graph import StateGraph, START, END
import os
import tempfile
import fitz
from pptx import Presentation
import subprocess
from io import BytesIO
from PIL import Image
import platform
import pandas as pd
import base64
import asyncio
from concurrent.futures import ThreadPoolExecutor
from openai import AsyncOpenAI

#my imports
from src.agent.graph_schemas import DocIngestionOutput, DocIngestionRequest, DataExtraction, IngestionState, FileData, SuccessStories
from src.agent.ingest import create_ingest_manager
from src.agent.logger import get_logger
from src.agent.config import Config
from src.agent.minio_config import MinIOManager
from src.agent.file_hash_manager import get_file_hash_manager

logger = get_logger(__name__)
config = Config()
# Global MinIO manager instance
minio_manager = MinIOManager()

# Global thread pool executor for better lifecycle management
_thread_pool_executor = ThreadPoolExecutor(max_workers=4)

try:
    import win32com.client
    WINDOWS_COM_AVAILABLE = True
except ImportError:
    WINDOWS_COM_AVAILABLE = False
    


# Docker-friendly path configuration for temporary operations only
WORK_DIR = config.WORK_DIR

logger.info(f"Work directory (temporary operations): {WORK_DIR}")
logger.info(f"Platform: {platform.system()}")
logger.info("File storage: MinIO (configured via MinIOManager)")
    
if not config.OPENROUTER_API_KEY:
    logger.warning("OPENROUTER_API_KEY not set in configuration.")


client = AsyncOpenAI(
    base_url=config.OPENAI_BASE_URL,
    api_key=config.OPENROUTER_API_KEY
)

# Helper functions for MinIO operations
def create_minio_object_path(filename: str, folder: str = "") -> str:
    """Create a MinIO object path with optional folder structure."""
    if folder:
        return f"{folder}/{filename}"
    return filename

def create_temp_file_from_minio(object_name: str, suffix: str = "") -> str:
    """Download a file from MinIO to a temporary location and return the path (synchronous - use create_temp_file_from_minio_async for better performance)."""
    data = minio_manager.get_object(object_name)
    if data is None:
        # Try to find the file with different possible paths
        logger.warning(f"Object {object_name} not found, trying alternative paths...")
        
        # Try just the filename if a full path was provided
        if "/" in object_name:
            filename_only = os.path.basename(object_name)
            logger.debug(f"Trying filename only: {filename_only}")
            data = minio_manager.get_object(filename_only)
        
        # Try with presentations/ prefix if just filename was provided
        if data is None and "/" not in object_name:
            presentations_path = f"presentations/{object_name}"
            logger.debug(f"Trying with presentations prefix: {presentations_path}")
            data = minio_manager.get_object(presentations_path)
        
        # List all available objects for debugging
        if data is None:
            available_objects = minio_manager.list_objects()
            logger.error(f"Object not found. Available objects: {available_objects}")
            raise FileNotFoundError(f"Object {object_name} not found in MinIO")
    
    # Create temporary file
    temp_fd, temp_path = tempfile.mkstemp(suffix=suffix)
    try:
        with os.fdopen(temp_fd, 'wb') as temp_file:
            temp_file.write(data)
        logger.debug(f"Created temporary file from MinIO object {object_name}: {temp_path}")
        return temp_path
    except Exception as e:
        os.close(temp_fd)
        if os.path.exists(temp_path):
            os.unlink(temp_path)
        raise e

async def create_temp_file_from_minio_async(object_name: str, suffix: str = "") -> str:
    """Download a file from MinIO to a temporary location and return the path (async)."""
    # Use the global thread pool for file operations
    loop = asyncio.get_event_loop()
    
    # Try to get the object asynchronously
    data = await minio_manager.get_object_async(object_name)
    if data is None:
        # Try to find the file with different possible paths
        logger.warning(f"Object {object_name} not found, trying alternative paths...")
        
        # Try just the filename if a full path was provided
        if "/" in object_name:
            filename_only = os.path.basename(object_name)
            logger.debug(f"Trying filename only: {filename_only}")
            data = await minio_manager.get_object_async(filename_only)
        
        # Try with presentations/ prefix if just filename was provided
        if data is None and "/" not in object_name:
            presentations_path = f"presentations/{object_name}"
            logger.debug(f"Trying with presentations prefix: {presentations_path}")
            data = await minio_manager.get_object_async(presentations_path)
        
        # List all available objects for debugging
        if data is None:
            available_objects = await minio_manager.list_objects_async()
            logger.error(f"Object not found. Available objects: {available_objects}")
            raise FileNotFoundError(f"Object {object_name} not found in MinIO")
    
    # Create temporary file using global thread pool
    def _create_temp_file():
        temp_fd, temp_path = tempfile.mkstemp(suffix=suffix)
        try:
            with os.fdopen(temp_fd, 'wb') as temp_file:
                temp_file.write(data)
            logger.debug(f"Created temporary file from MinIO object {object_name}: {temp_path}")
            return temp_path
        except Exception as e:
            os.close(temp_fd)
            if os.path.exists(temp_path):
                os.unlink(temp_path)
            raise e
    
    return await loop.run_in_executor(_thread_pool_executor, _create_temp_file)

def cleanup_temp_file(temp_path: str) -> None:
    """Safely remove a temporary file (synchronous)."""
    try:
        if os.path.exists(temp_path):
            os.unlink(temp_path)
            logger.debug(f"Cleaned up temporary file: {temp_path}")
    except Exception as e:
        logger.warning(f"Failed to cleanup temporary file {temp_path}: {e}")

async def cleanup_temp_file_async(temp_path: str) -> None:
    """Safely remove a temporary file (async)."""
    def _cleanup():
        try:
            if os.path.exists(temp_path):
                os.unlink(temp_path)
                logger.debug(f"Cleaned up temporary file: {temp_path}")
        except Exception as e:
            logger.warning(f"Failed to cleanup temporary file {temp_path}: {e}")
    
    loop = asyncio.get_event_loop()
    await loop.run_in_executor(_thread_pool_executor, _cleanup)

async def store_file_to_minio(file_path: str, object_name: str, content_type: str = "application/octet-stream") -> bool:
    """Store a local file to MinIO and return success status."""
    def _read_file():
        with open(file_path, 'rb') as file_data:
            return file_data.read()
    
    try:
        loop = asyncio.get_event_loop()
        content = await loop.run_in_executor(_thread_pool_executor, _read_file)
        return await minio_manager.put_object_async(object_name, content, content_type)
    except Exception as e:
        logger.error(f"Failed to store file {file_path} to MinIO as {object_name}: {e}")
        return False
def get_safe_file_path(file_path: str) -> str:
    """Convert file path to be container-safe and absolute."""
    # If it's a MinIO object path (contains no path separators or starts with folder structure)
    if '/' in file_path and not os.path.isabs(file_path):
        return file_path  # Return as-is for MinIO object paths
    if os.path.isabs(file_path):
        return file_path
    return os.path.join(WORK_DIR, file_path)


async def process_uploaded_files(state: IngestionState):
    """Process uploaded files from FileData objects and save them to MinIO with duplicate detection."""
    if not state.files or len(state.files) == 0:
        logger.warning("No files provided for processing")
        return {"status": "No files to process"}

    logger.info(f"Processing {len(state.files)} uploaded files with duplicate detection")

    # Get the file hash manager
    hash_manager = get_file_hash_manager()

    for file_obj in state.files:
        filename_lower = file_obj.filename.lower()

        # Decode the base64 content first (needed for both hashing and storage)
        try:
            file_content = base64.b64decode(file_obj.content)
        except Exception as e:
            logger.error(f"Error decoding base64 content for {file_obj.filename}: {e}")
            raise Exception(f"Failed to decode file content for {file_obj.filename}")

        # Calculate file hash (use provided hash if available, otherwise calculate)
        if file_obj.file_hash:
            file_hash = file_obj.file_hash
            logger.debug(f"Using provided hash for {file_obj.filename}: {file_hash[:16]}...")
        else:
            file_hash = await hash_manager.calculate_file_hash(file_content)
            logger.debug(f"Calculated hash for {file_obj.filename}: {file_hash[:16]}...")

        # Check for duplicates
        is_duplicate, existing_metadata = await hash_manager.check_duplicate(file_hash)

        if is_duplicate:
            # File is a duplicate - return error with detailed information
            error_msg = (
                f"Duplicate file detected: '{file_obj.filename}' has already been uploaded.\n"
                f"Original file: '{existing_metadata.filename}'\n"
                f"Original upload date: {existing_metadata.upload_timestamp}\n"
                f"Original location: {existing_metadata.minio_object_path}\n"
                f"File hash: {file_hash[:16]}..."
            )
            logger.warning(error_msg)
            raise Exception(error_msg)

        # File is not a duplicate - proceed with storage based on file type
        if filename_lower.endswith('.pptx'):
            object_name = create_minio_object_path(file_obj.filename, "presentations")

            try:
                success = await minio_manager.put_object_async(
                    object_name,
                    file_content,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

                if success:
                    # Register the file in hash registry
                    await hash_manager.register_file(
                        file_hash=file_hash,
                        filename=file_obj.filename,
                        file_size=len(file_content),
                        content_type=file_obj.content_type,
                        minio_object_path=object_name
                    )
                    logger.info(f"Saved PowerPoint file to MinIO: {object_name}")
                    state.ppt_path = object_name
                    break
                else:
                    raise Exception(f"Failed to store file {file_obj.filename} to MinIO")

            except Exception as e:
                logger.error(f"Error saving file {file_obj.filename} to MinIO: {e}")
                raise e

        elif filename_lower.endswith('.pdf'):
            object_name = create_minio_object_path(file_obj.filename, "pdfs")

            try:
                success = await minio_manager.put_object_async(
                    object_name,
                    file_content,
                    "application/pdf"
                )

                if success:
                    # Register the file in hash registry
                    await hash_manager.register_file(
                        file_hash=file_hash,
                        filename=file_obj.filename,
                        file_size=len(file_content),
                        content_type=file_obj.content_type,
                        minio_object_path=object_name
                    )
                    logger.info(f"Saved PDF file to MinIO: {object_name}")
                    state.pdf_path = object_name
                    break
                else:
                    raise Exception(f"Failed to store file {file_obj.filename} to MinIO")

            except Exception as e:
                logger.error(f"Error saving file {file_obj.filename} to MinIO: {e}")
                raise e

        elif filename_lower.endswith(('.docx', '.doc')):
            object_name = create_minio_object_path(file_obj.filename, "documents")

            try:
                content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document" if filename_lower.endswith('.docx') else "application/msword"

                success = await minio_manager.put_object_async(
                    object_name,
                    file_content,
                    content_type
                )

                if success:
                    # Register the file in hash registry
                    await hash_manager.register_file(
                        file_hash=file_hash,
                        filename=file_obj.filename,
                        file_size=len(file_content),
                        content_type=file_obj.content_type,
                        minio_object_path=object_name
                    )
                    logger.info(f"Saved Word document to MinIO: {object_name}")
                    state.doc_path = object_name
                    break
                else:
                    raise Exception(f"Failed to store file {file_obj.filename} to MinIO")

            except Exception as e:
                logger.error(f"Error saving file {file_obj.filename} to MinIO: {e}")
                raise e

    # Check what type of file was processed and return appropriate status
    
    if state.ppt_path:
        source_document_link = state.ppt_path
        return {"source_document_link": source_document_link, "ppt_path": state.ppt_path, "status": "PowerPoint file processed successfully (no duplicates)"}
    elif state.pdf_path:
        source_document_link = state.pdf_path
        return {"source_document_link": source_document_link, "pdf_path": state.pdf_path, "status": "PDF file processed successfully (no duplicates)"}
    elif state.doc_path:
        source_document_link = state.doc_path
        return {"source_document_link": source_document_link, "doc_path": state.doc_path, "status": "Word document processed successfully (no duplicates)"}
    else:
        logger.warning("No supported file found in uploaded files (supported: .pptx, .pdf, .docx, .doc)")
        return {"status": "No supported file found"}

    return {"status": "Files processed successfully"}

async def convert_to_pdf_node(state: IngestionState):
    """Convert a PowerPoint or Word document to PDF using the best available method (Windows COM or LibreOffice)."""
    # Determine which type of file we're converting
    source_path = None
    source_type = None
    file_suffix = None
    folder_prefix = None
    
    if state.ppt_path:
        source_path = state.ppt_path
        source_type = "PowerPoint"
        file_suffix = ".pptx"
        folder_prefix = "presentations/"
    elif state.doc_path:
        source_path = state.doc_path
        source_type = "Word"
        file_suffix = ".docx" if state.doc_path.lower().endswith('.docx') else ".doc"
        folder_prefix = "documents/"
    else:
        logger.info("No PowerPoint or Word document found in state")
        return {"status": "No PowerPoint or Word document found"}
    
    logger.info(f"Starting {source_type} to PDF conversion for MinIO object: {source_path}")
    logger.debug(f"Received {source_type.lower()}_path from state: '{source_path}'")

    # Ensure source_path has the correct folder prefix
    if source_path and not source_path.startswith(folder_prefix):
        logger.info(f"Adding {folder_prefix} prefix to {source_type.lower()}_path: {source_path}")
        source_path = f"{folder_prefix}{source_path}"
        # Update the state with the corrected path
        if source_type == "PowerPoint":
            state.ppt_path = source_path
        else:
            state.doc_path = source_path

    # List available objects to debug
    available_objects = await minio_manager.list_objects_async(prefix=folder_prefix)
    logger.debug(f"Available objects in {folder_prefix}: {available_objects}")

    # Create temporary file for document processing
    temp_source_path = None
    temp_pdf_path = None
    
    try:
        # Download source file from MinIO to temporary file (using async version)
        temp_source_path = await create_temp_file_from_minio_async(source_path, suffix=file_suffix)
        
        # Create temporary PDF file
        def _create_temp_pdf():
            temp_pdf_fd, temp_pdf_path = tempfile.mkstemp(suffix=".pdf")
            os.close(temp_pdf_fd)  # Close the file descriptor, we'll use the path
            return temp_pdf_path
        
        loop = asyncio.get_event_loop()
        temp_pdf_path = await loop.run_in_executor(_thread_pool_executor, _create_temp_pdf)
        
        base_name = os.path.basename(source_path)
        source_name_without_ext = os.path.splitext(base_name)[0].replace(' ', '_')
        pdf_object_name = create_minio_object_path(f"{source_name_without_ext}.pdf", "pdfs")
        
        logger.debug(f"Target PDF MinIO object: {pdf_object_name}")

        try:
            logger.info(f"Using LibreOffice conversion for {source_type}...")
            # Use absolute paths and ensure proper quoting for Docker/Linux environments
            abs_source_path = os.path.abspath(temp_source_path)
            temp_output_dir = os.path.dirname(temp_pdf_path)
                
            command = [
                'libreoffice', 
                '--headless', 
                '--convert-to', 'pdf', 
                '--outdir', temp_output_dir, 
                abs_source_path
            ]
                
            logger.debug(f"LibreOffice command: {' '.join(command)}")
            
            # Run LibreOffice in thread pool to avoid blocking
            def _run_libreoffice():
                return subprocess.run(command, check=True, capture_output=True, text=True, timeout=300)
            
            result = await loop.run_in_executor(_thread_pool_executor, _run_libreoffice)
            logger.debug(f"LibreOffice stdout: {result.stdout}")
                
            # LibreOffice creates PDF with same name as source, find the generated file
            expected_pdf = os.path.join(temp_output_dir, f"{os.path.splitext(os.path.basename(temp_source_path))[0]}.pdf")
            if os.path.exists(expected_pdf) and expected_pdf != temp_pdf_path:
                # Move to our target temp path
                def _move_file():
                    os.rename(expected_pdf, temp_pdf_path)
                
                await loop.run_in_executor(_thread_pool_executor, _move_file)
                
            logger.info(f"LibreOffice {source_type} conversion successful")
        except subprocess.TimeoutExpired as e:
            logger.error(f"LibreOffice conversion timed out after 5 minutes: {e}")
            raise e
        except subprocess.CalledProcessError as e:
            logger.error(f"LibreOffice conversion failed with exit code {e.returncode}: {e.stderr}")
            raise e
        except Exception as e:
            logger.error(f"LibreOffice conversion failed: {e}")
            raise e
        

        success = await store_file_to_minio(temp_pdf_path, pdf_object_name, "application/pdf")
        if not success:
            raise Exception(f"Failed to store PDF to MinIO: {pdf_object_name}")
        
        logger.info(f"Successfully stored PDF to MinIO: {pdf_object_name}")
        state.pdf_path = pdf_object_name
        return {"pdf_path": state.pdf_path, "status": f"{source_type} converted to PDF successfully"}
        
    finally:
        if temp_source_path:
            await cleanup_temp_file_async(temp_source_path)
        if temp_pdf_path:
            await cleanup_temp_file_async(temp_pdf_path)
        
async def convert_pdf_to_images_node(state: IngestionState):
    """Convert a PDF file to a series of images using PyMuPDF."""
    pdf_object_name = state.pdf_path
    
    # Check if PDF path is available
    if not pdf_object_name or pdf_object_name.strip() == "":
        logger.warning("No PDF path available for conversion to images")
        return {"status": "No PDF file to convert to images"}
    
    logger.info(f"Starting PDF to images conversion for MinIO object: {pdf_object_name}")

    temp_pdf_path = None
    
    try:
        # Download PDF from MinIO to temporary file (using async version)
        temp_pdf_path = await create_temp_file_from_minio_async(pdf_object_name, suffix=".pdf")
        
        # Open PDF and convert to images
        def _process_pdf():
            doc = fitz.open(temp_pdf_path)
            base_name = os.path.basename(pdf_object_name)
            pdf_name_without_ext = os.path.splitext(base_name)[0].replace(' ', '_')
            pages_data = []
            
            logger.debug(f"Converting {len(doc)} pages to images")
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap()  # Convert the page to a pixmap (image)
                
                # Create image data in memory
                image_data = pix.tobytes("png")
                
                # Create MinIO object name for the image
                image_object_name = create_minio_object_path(
                    f"{pdf_name_without_ext}_{page_num:04d}.png", 
                    "images"
                )
                
                pages_data.append((image_object_name, page_num, image_data))
            
            doc.close()
            return pages_data, pdf_name_without_ext
        
        loop = asyncio.get_event_loop()
        pages_data, pdf_name_without_ext = await loop.run_in_executor(_thread_pool_executor, _process_pdf)
        
        # Store images to MinIO asynchronously
        image_paths = []
        for image_object_name, page_num, image_data in pages_data:
            success = await minio_manager.put_object_async(image_object_name, image_data, "image/png")
            if success:
                image_paths.append((image_object_name, page_num))
                logger.debug(f"Stored image to MinIO: {image_object_name}")
            else:
                logger.warning(f"Failed to store image to MinIO: {image_object_name}")
        
        logger.info(f"Successfully converted {len(image_paths)} PDF pages to images in MinIO")
        state.image_paths_list = image_paths
        return {"image_paths_list": state.image_paths_list, "status": "PDF converted to images successfully"}
        
    finally:
        # Cleanup temporary PDF file (using async version)
        if temp_pdf_path:
            await cleanup_temp_file_async(temp_pdf_path)

def get_b64_image_from_content(image_content):
    """Convert image content to a base64 encoded string."""
    img = Image.open(BytesIO(image_content))
    if img.mode != 'RGB':
        img = img.convert('RGB')
    buffered = BytesIO()
    img.save(buffered, format="JPEG")

    return base64.b64encode(buffered.getvalue()).decode("utf-8")

async def get_b64_image_from_content_async(image_content):
    """Convert image content to a base64 encoded string (async)."""
    def _convert_image():
        img = Image.open(BytesIO(image_content))
        if img.mode != 'RGB':
            img = img.convert('RGB')
        buffered = BytesIO()
        img.save(buffered, format="JPEG")
        return base64.b64encode(buffered.getvalue()).decode("utf-8")
    
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(_thread_pool_executor, _convert_image)

async def describe_image(image_content, previous_extractions=None):
    """Generate a description of an image."""
    logger.debug("Starting image description using OpenAI API")
    
    # Convert the image content to a base64 encoded string (using async version)
    image_b64 = await get_b64_image_from_content_async(image_content)
    
    # Prepare context from previous extractions if available
    context_section = ""
    if previous_extractions and len(previous_extractions) > 0:
        context_section = "\n\nPREVIOUS SLIDE CONTEXT (for understanding story progression - DO NOT directly use this data):\n"
        for i, prev_data in enumerate(previous_extractions, 1):
            context_section += f"\n--- Previous Slide {i} Summary ---\n"
            for key, value in prev_data.items():
                if value and str(value).strip() and str(value) != 'nan' and key not in ['source_document_link', 'slide_number']:
                    context_section += f"{key}: {str(value)[:200]}{'...' if len(str(value)) > 200 else ''}\n"
        context_section += "\nIMPORTANT: Use this context only to understand the story flow\n"
    
    response = await client.beta.chat.completions.parse(
        model=config.LLM_MODEL,

        messages=[
            {
                "role": "system",
                "content": """You are an expert at extracting comprehensive technical success story data for Reailize's Sales Enablement Chat System. 
                
Your goal is to extract ALL available data that will be used by sales teams to build credibility with potential clients and accelerate sales cycles by providing evidence-based value statements and compelling success stories.

Focus on capturing COMPLETE and DETAILED information without summarizing or condensing. Extract information EXACTLY as stated in the slide content and preserve ALL technical specifications, business value, and project details."""
            },
            {
                "role": "user", 
                "content": [{
                    "type": "text",
                    "text": f"""Extract ALL success stories from this PowerPoint slide for the Sales Enablement Knowledge Base with MAXIMUM DETAIL AND COMPLETENESS.

CONTEXT: This slide contains delivery experiences that need to be captured in their entirety as sales-oriented success stories. Extract ALL available information for:

STORY TYPES TO CAPTURE:
- Customer implementation projects and technical solutions
- Product success stories and feature developments
- Internal tool developments and process improvements
- Partnership collaborations and joint solutions

SUBJECT & CONTEXT - CAPTURE ALL DETAILS:
- Subject name (customer organization, product name, or project subject - exact name as stated)
- Story type classification (customer, product, internal, or partnership)
- ALL business challenges, operational issues, strategic pain points, and problems before intervention
- ALL technical challenges, system limitations, infrastructure constraints, and technology gaps
- Include ALL context about industry, business environment, regulatory requirements, or operational constraints

SOLUTION ARCHITECTURE & DELIVERY - COMPLETE TECHNICAL DETAILS:
- FULL description of solution boundaries, scale, scope, and extent (include everything that was and wasn't included)
- ALL concrete deliverables, implementations, outcomes, and work products provided
- ALL key capabilities, functionalities, and core competencies that the solution provides or enables
- ALL specific features, components, modules, or functionality elements that comprise the solution
- COMPLETE list of technologies, tools, platforms, software, hardware, and technical components used
- ALL technical competencies, expertise areas, specialized skills, and capabilities demonstrated
- ALL network infrastructure components, telecommunications elements, vendor partnerships, and third-party integrations
- Include ALL technical architectures, system designs, integration points, and implementation approaches

PROJECT EXECUTION & CHALLENGES - COMPREHENSIVE PROJECT DETAILS:
- ALL significant problems, obstacles, complications encountered and DETAILED descriptions of how they were resolved
- COMPLETE timeline from project/product initiation to completion, including ALL milestones, phases, and key dates
- EXACT number of team members involved, ALL roles, expertise areas, responsibilities, and team composition
- ALL changes in project scope, timeline, team composition, requirements, and evolution throughout implementation
- COMPLETE description of collaboration patterns, working relationships, communication approaches, coordination methods, and team dynamics
- Include ALL project management approaches, methodologies, tools, and processes used

BUSINESS VALUE & OUTCOMES - ALL MEASURABLE RESULTS:
- ALL measurable business outcomes and quantifiable value delivered (revenue increase, cost reduction, ROI, time savings, performance improvements)
- ALL overall benefits achieved including qualitative improvements, strategic advantages, competitive benefits, user satisfaction, and long-term value creation
- COMPLETE list of key stakeholders, decision makers, sponsors, technical contacts, and anyone who can serve as references
- ALL evidence of success, customer satisfaction, feedback, testimonials, and post-implementation results
- Include ALL metrics, KPIs, performance indicators, and measurable outcomes achieved

KEYWORDS & TERMINOLOGY EXTRACTION - SEARCHABILITY & CATEGORIZATION:
- Extract ALL key terms, concepts, technologies, important phrases, and technical terminology for searchability and categorization
- Create a comprehensive glossary of ALL acronyms, technical terms, specialized terminology, and domain-specific language with their definitions
- Include ALL technology names, product names, vendor names, industry terms, and specialized concepts
- Capture ALL relevant search terms that would help find this success story in future queries

EXTRACTION INSTRUCTIONS - MAXIMUM DETAIL CAPTURE:
1. Each distinct project, product development, or use case should be a separate story entry with FULL details
2. Extract information EXACTLY as stated in the slide - include ALL available details without condensing
3. Preserve ALL technical specifications, product names, version numbers, quantities, and technical parameters
4. Include ALL metrics, timelines, quantifiable outcomes, success indicators, and performance data
5. Capture ALL vendor names, technology brands, product names, and technical ecosystem details
6. Include ALL project specifics: dates, durations, team sizes, budget information (if mentioned), and scope details
7. Preserve ALL customer quotes, feedback, testimonials, or satisfaction indicators
8. Extract ALL keywords, key terms, concepts, technologies, and important phrases for searchability
9. Create a comprehensive glossary mapping ALL acronyms and technical terms to their definitions
10. If no clear success story is present in the slide, return an empty stories list
11. DO NOT summarize, condense, or shorten any information - capture everything in full detail

{context_section}
                        """},
                            {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/png;base64,{image_b64}"
                                    }
                            }
                        ],
            }
        ],
        response_format=DataExtraction
    )

    return response

async def extract_data_node(state: IngestionState):
    """Process a PowerPoint file."""
    logger.info("Starting data extraction from PowerPoint slides")

    # Ensure we have valid lists, not None
    #slide_texts = state.text_and_notes
    images_data = state.image_paths_list
    
    # if slide_texts is None:
    #     slide_texts = []
    #     logger.warning("text_and_notes is None, using empty list")
    
    if images_data is None:
        images_data = []
        logger.warning("image_paths_list is None, using empty list")
    
    if not images_data or len(images_data) == 0:
        logger.warning("No images available for data extraction - skipping extraction")
        return {"ExtractedData_list": [], "status": "No images to extract data from"}
    
    ExtractedData_list = []
    
    logger.info(f"Processing slides: {len(images_data)} images")
    
    # Ensure we process all images, even if there are more images than text slides
    if len(images_data) == 0:
        logger.warning("No images to process")
        return {"ExtractedData_list": ExtractedData_list, "status": "No images to process"}

    max_items = len(images_data)

    for i in range(max_items):
        # Get image data if available
        if i < len(images_data):
            image_object_name, page_num = images_data[i]
            logger.info(f"Processing slide {page_num + 1} (index {i})")
        else:
            logger.warning(f"No image data for index {i}")
            continue
        
        # # Get text data if available, otherwise use empty text
        # if i < len(slide_texts):
        #     text_note_dict = slide_texts[i]
        #     slide_text = text_note_dict["text"]
        #     notes = text_note_dict["notes"]
        # else:
        #     logger.info(f"No text data for slide {page_num + 1}, using empty text")
        #     slide_text = ""
        #     notes = ""
        
        # if notes:
        #     notes = "\n\nThe speaker notes for this slide are: " + notes
        
        try:
            # Get image content from MinIO (using async version)
            image_content = await minio_manager.get_object_async(image_object_name)
            if image_content is None:
                logger.error(f"Failed to retrieve image from MinIO: {image_object_name}")
                continue

            response = await describe_image(image_content, ExtractedData_list)
            extracted_data = response.choices[0].message.parsed
            
            if extracted_data and extracted_data.ExtractedData:
                logger.debug(f"Found {len(extracted_data.ExtractedData)} stories in slide {page_num + 1}")
                for data in extracted_data.ExtractedData:
                    data_dict = data.model_dump()
                    for key, value in data_dict.items():
                        if isinstance(value, dict) and value:
                            # Handle glossary dictionary - convert to JSON string for storage
                            import json
                            data_dict[key] = json.dumps(value, ensure_ascii=False)
                        elif isinstance(value, dict):
                            data_dict[key] = ""
                        elif isinstance(value, list) and value:
                            data_dict[key] = "; ".join(str(item) for item in value if item)
                        elif isinstance(value, list):
                            data_dict[key] = ""
                        elif value is None:
                            data_dict[key] = ""
                    
                    data_dict["source_document_link"] = state.source_document_link
                    data_dict["slide_number"] = page_num + 1

                    ExtractedData_list.append(data_dict)
                    logger.info(f"extracted data:   {data_dict}")
            else:
                logger.debug(f"No stories found in slide {page_num + 1}")
        except Exception as e:
            logger.error(f"Error processing slide {page_num + 1}: {e}")
            continue
            
    if ExtractedData_list:
        df = pd.DataFrame(ExtractedData_list)
        logger.info(f"Extracted {len(ExtractedData_list)} success stories for Sales Enablement Knowledge Base")
        
        # Create a comprehensive output file with timestamp
        import datetime
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # Save CSV to MinIO instead of local filesystem
        csv_object_name = create_minio_object_path(f"success_stories_{timestamp}.csv", "reports")
        latest_csv_object_name = create_minio_object_path("success_stories_latest.csv", "reports")
        
        # Convert DataFrame to CSV bytes
        csv_buffer = BytesIO()
        df.to_csv(csv_buffer, index=False)
        csv_content = csv_buffer.getvalue()
        
        # Store both timestamped and latest versions to MinIO
        success1 = await minio_manager.put_object_async(csv_object_name, csv_content, "text/csv")
        success2 = await minio_manager.put_object_async(latest_csv_object_name, csv_content, "text/csv")
        
        if success1 and success2:
            logger.info(f"Success stories saved to MinIO: {csv_object_name}")
            logger.info(f"Latest version saved to MinIO: {latest_csv_object_name}")
        else:
            logger.warning("Failed to save some CSV files to MinIO")
        
        state.ExtractedData_list = ExtractedData_list
        return {"ExtractedData_list": state.ExtractedData_list, "status": "Data extraction completed successfully"}
    else:
        logger.warning("No success stories found in the document")
        state.ExtractedData_list = []
        return {"ExtractedData_list": state.ExtractedData_list, "status": "No success stories found"}



# def extract_text_and_notes_from_ppt_node(state: IngestionState):
#     """Extract text and notes from a PowerPoint file."""
#     ppt_object_name = state.ppt_path
#     logger.info(f"Starting text and notes extraction from MinIO PowerPoint object: {ppt_object_name}")

#     # Ensure ppt_object_name has the correct presentations/ prefix
#     if ppt_object_name and not ppt_object_name.startswith("presentations/"):
#         logger.info(f"Adding presentations/ prefix to ppt_path: {ppt_object_name}")
#         ppt_object_name = f"presentations/{ppt_object_name}"
#         # Update the state with the corrected path
#         state.ppt_path = ppt_object_name

#     temp_ppt_path = None
    
#     try:
#         # Download PPT from MinIO to temporary file
#         temp_ppt_path = create_temp_file_from_minio(ppt_object_name, suffix=".pptx")
        
#         # Extract text and notes using python-pptx
#         prs = Presentation(temp_ppt_path)
#         text_and_notes = []
        
#         logger.debug(f"Processing {len(prs.slides)} slides for text extraction")
        
#         for slide in prs.slides:
#             slide_text = ' '.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
#             try:
#                 notes = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ''
#             except Exception as e:
#                 logger.debug(f"Could not extract notes from slide: {e}")
#                 notes = ''
            
#             text_and_notes.append({"text": slide_text, "notes": notes})
        
#         logger.info(f"Successfully extracted text and notes from {len(text_and_notes)} slides")
#         state.text_and_notes = text_and_notes
#         return {"text_and_notes": state.text_and_notes, "status": "Text and notes extraction completed successfully"}
        
    # finally:
    #     # Cleanup temporary PPT file
    #     if temp_ppt_path:
    #         cleanup_temp_file(temp_ppt_path)
    
async def star_framework_node(state: IngestionState):
    """Generate a STAR framework story based on extracted success stories."""
    logger.info("Starting STAR framework story generation")
    
    # Get the extracted success stories from state
    stories_data = state.ExtractedData_list

    if (stories_data is None or len(stories_data) == 0) and not state.additional_context:
        logger.warning("No success stories available for STAR framework generation")
        return {"status": "No success stories found to generate a STAR framework story.", "star_story": ""}
    
    stories_context = stories_data
    
    # Create a formatted context string
    context_text = ""
    
    # Add user input as email content if available
    if state.additional_context and state.additional_context.strip():
        context_text += f"\n--- Email Content ---\n"
        context_text += f"User Input: {state.additional_context}\n"
    
    for i, story in enumerate(stories_context, 1):
        context_text += f"\n--- Extracted Data {i} ---\n"
        for key, value in story.items():
            if value and str(value).strip() and str(value) != 'nan':
                context_text += f"{key}: {value}\n"
    
    logger.info(f"Generated context for STAR framework:\n{context_text}")
    
    response = await client.beta.chat.completions.parse(
        model=config.LLM_MODEL,
        messages=[
            {
                "role": "system",
                "content": """You are an expert storyteller specializing in business case studies for Reailize.
                            You will structure responses using the STAR Framework:
                            1) The Star (the client/hero) - Who is the customer and what challenges did they face?
                            2) The Sage (our team/solution) - How did Reailize guide/enable them with our expertise and solutions?
                            3) The Shift (the transformation/outcome) - What tangible positive outcomes were achieved?
                            
                            Your role is to create clear, aspirational, and outcome-oriented stories that demonstrate Reailize's value proposition and technical capabilities."""
            },
            {
                "role": "user",
                "content": f"""Based on the following extracted data, write a compelling case study story using the STAR Framework.

EXTRACTED DATA:
{context_text}

Please create a narrative that:
- Identifies the Star (client/customer) and their specific challenges
- Describes the Sage (Reailize's role) and how we provided solutions
- Highlights the Shift (transformation/outcomes) with quantifiable results when available

Structure your response as:
**THE STAR:** [Customer and their challenges]
**THE SAGE:** [Reailize's solution and approach]
**THE SHIFT:** [Outcomes and transformation achieved]

Make it clear, concise, client-focused, and suitable for sales enablement."""
            }
        ],
    )
    
    star_story = response.choices[0].message.content
    logger.info("STAR framework story generation completed")
    return {"status": "STAR framework story generated successfully", "star_story": star_story}


async def hero_framework_node(state: IngestionState):
    """Generate a HERO framework story based on extracted success stories."""
    logger.info("Starting HERO framework story generation")

    # Get the extracted success stories from state
    stories_data = state.ExtractedData_list
    
    if (stories_data is None or len(stories_data) == 0) and not state.additional_context:
        logger.warning("No success stories available for hero framework generation")
        return {"status": "No success stories found to generate a hero framework story.", "hero_story": ""}
    
    stories_context = stories_data
    
    # Create a formatted context string
    context_text = ""
    
    if state.additional_context and state.additional_context.strip():
        context_text += f"\n--- Email Content ---\n"
        context_text += f"User Input: {state.additional_context}\n"
    
    for i, story in enumerate(stories_context, 1):
        context_text += f"\n--- Extracted Data {i} ---\n"
        for key, value in story.items():
            if value and str(value).strip() and str(value) != 'nan':
                context_text += f"{key}: {value}\n"
    
    logger.info(f"Generated context for hero framework:\n{context_text}")
    
    response = await client.beta.chat.completions.parse(
        model=config.LLM_MODEL,
        messages=[
            {
                "role": "system",
                "content": """You are an expert storyteller specializing in business case studies for Reailize.
                            You will structure responses using the Hero–Villain–Guide Framework:
                            1) The Hero: Who is the client and what are they trying to achieve?
                            2) The Villain: What obstacle/problem stood in their way?
                            3) The Guide: How did we (our team, company, or technology) provide the plan, tools, or expertise to help them win?
                            
                            Your role is to create clear, aspirational, and outcome-oriented stories that demonstrate Reailize's value proposition and technical capabilities."""
            },
            {
                "role": "user",
                "content": f"""Based on the following extracted data, write a compelling case study story using the Hero–Villain–Guide Framework.

EXTRACTED DATA:
{context_text}

Please create a narrative that:
- Identifies the Hero (client/customer) and what they were trying to achieve
- Describes the Villain (obstacles/problems) that stood in their way
- Explains how the Guide (Reailize) provided the plan, tools, or expertise to help them win

Structure your response as:
**THE HERO:** [Customer and their goals/aspirations]
**THE VILLAIN:** [Obstacles, challenges, and problems they faced]
**THE GUIDE:** [How Reailize provided solutions, expertise, and guidance to help them succeed]

Make it clear, concise, client-focused, and suitable for sales enablement."""
            }
        ],
    )
    
    hero_story = response.choices[0].message.content
    logger.info("hero framework story generation completed")
    return {"status": "Hero framework story generated successfully", "hero_story": hero_story} 


async def pas_framework_node(state: IngestionState):
    """Generate a Problem-Agitate-Solve framework story based on extracted success stories."""
    logger.info("Starting PAS framework story generation")

    # Get the extracted success stories from state
    stories_data = state.ExtractedData_list
    
    if (stories_data is None or len(stories_data) == 0) and not state.additional_context:
        logger.warning("No success stories available for PAS framework generation")
        return {"status": "No success stories found to generate a PAS framework story.", "pas_story": ""}
    
    stories_context = stories_data
    
    # Create a formatted context string
    context_text = ""
    
    if state.additional_context and state.additional_context.strip():
        context_text += f"\n--- Email Content ---\n"
        context_text += f"User Input: {state.additional_context}\n"
    
    for i, story in enumerate(stories_context, 1):
        context_text += f"\n--- Extracted Data {i} ---\n"
        for key, value in story.items():
            if value and str(value).strip() and str(value) != 'nan':
                context_text += f"{key}: {value}\n"
    
    logger.info(f"Generated context for PAS framework:\n{context_text}")
    
    response = await client.beta.chat.completions.parse(
        model=config.LLM_MODEL,
        messages=[
            {
                "role": "system",
                "content": """You are an expert storyteller specializing in business case studies for Reailize.
                            You will structure responses using the Problem-Agitate-Solve Framework:
                            1) The Problem: What is the main issue or challenge the client is facing?
                            2) The Agitation: Why is this problem significant, and what are the consequences of not addressing it? (Explain the wasted time, risks, costs, or missed opportunities)
                            3) The Solution: How did we (our team, company, or technology) provide the plan, tools, or expertise to help them overcome the challenges?

                            Your role is to create clear, aspirational, and outcome-oriented stories that demonstrate Reailize's value proposition and technical capabilities."""
            },
            {
                "role": "user",
                "content": f"""Based on the following extracted data, write a compelling case study story using the Problem-Agitate-Solve Framework.

EXTRACTED DATA:
{context_text}

Please create a narrative that:
- Identifies the Problem (main issues/challenges the client was facing)
- Agitates the Problem (explains why this was significant and the consequences of not addressing it)
- Presents the Solution (how Reailize provided the plan, tools, or expertise to solve the problems)

Structure your response as:
**THE PROBLEM:** [Main issues, challenges, and pain points the customer was experiencing]
**THE AGITATION:** [Why these problems were critical - costs, risks, missed opportunities, and consequences of inaction]
**THE SOLUTION:** [How Reailize provided the expertise, tools, and solutions to resolve the problems and deliver value]

Make it clear, concise, client-focused, and suitable for sales enablement."""
            }
        ],
    )
    
    pas_story = response.choices[0].message.content
    logger.info("PAS framework story generation completed")
    return {"status": "PAS framework story generated successfully", "pas_story": pas_story} 


async def validate_combined_story_node(state: IngestionState):
    """Validate and generate a unified comprehensive success story combining insights from all framework stories."""
    logger.info("Starting combined story validation and generation")

    # Get the extracted success stories from state
    stories_data = state.ExtractedData_list
    star_story = state.star_story
    hero_story = state.hero_story
    pas_story = state.pas_story
    
    if (stories_data is None or len(stories_data) == 0) and not state.additional_context:
        logger.warning("No success stories available for combined story generation")
        return {"status": "No success stories found to generate a combined story.", "combined_story": ""}
    
    if not any([star_story, hero_story, pas_story]):
        logger.warning("No framework stories available for combined story generation")
        return {"status": "No framework stories found to generate a combined story.", "combined_story": ""}
    

    context_text = ""

    if state.additional_context and state.additional_context.strip():
        context_text += f"\n--- Email Content ---\n"
        context_text += f"User Input: {state.additional_context}\n"
    
    for i, story in enumerate(stories_data, 1):
        context_text += f"\n--- Extracted Data {i} ---\n"
        for key, value in story.items():
            if value and str(value).strip() and str(value) != 'nan':
                context_text += f"{key}: {value}\n"
    
    # Combine all framework stories for validation
    framework_stories_text = ""
    available_frameworks = []
    if star_story and star_story.strip():
        framework_stories_text += f"STAR Framework Story:\n{star_story}\n\n"
        available_frameworks.append("STAR")
    if hero_story and hero_story.strip():
        framework_stories_text += f"HERO Framework Story:\n{hero_story}\n\n"
        available_frameworks.append("HERO")
    if pas_story and pas_story.strip():
        framework_stories_text += f"PAS Framework Story:\n{pas_story}\n\n"
        available_frameworks.append("PAS")
    
    logger.info(f"Available frameworks for validation and combination: {available_frameworks}")
    
    # First, validate the framework stories
    validation_response = await client.beta.chat.completions.parse(
        model=config.LLM_MODEL,
        messages=[
            {
                "role": "system",
                "content": """You are an expert quality assurance analyst for business case studies and success stories.
                            Your task is to validate framework stories (STAR, HERO, PAS) for:
                            
                            1. COMPLETENESS: Does each story contain all necessary elements for its framework?
                            2. CONSISTENCY: Are the facts and details consistent across all framework stories?
                            3. CLARITY: Are the stories clear, well-structured, and easy to understand?
                            4. BUSINESS VALUE: Do the stories effectively communicate measurable business outcomes?
                            5. ACCURACY: Are there any obvious contradictions or missing critical information?
                            
                            Provide specific feedback on what works well and what needs improvement in each story.
                            Identify any gaps, inconsistencies, or areas that could be strengthened."""
            },
            {
                "role": "user",
                "content": f"""Please validate the following framework stories for quality, completeness, and consistency:

ORIGINAL EXTRACTED DATA (for reference):
{context_text}

FRAMEWORK STORIES TO VALIDATE:
{framework_stories_text}

Provide a validation report that includes:
1. Overall quality assessment (Good/Needs Improvement/Poor)
2. Specific strengths of each framework story
3. Identified gaps, inconsistencies, or areas for improvement
4. Recommendations for creating a comprehensive unified story
5. Key elements that should be emphasized in the final combined story

Be specific and constructive in your feedback."""
            }
        ],
    )
    
    validation_feedback = validation_response.choices[0].message.content
    logger.info("Framework stories validation completed")
    logger.debug(f"Validation feedback: {validation_feedback}")
    
    # Now create unique success stories incorporating validation insights
    response = await client.beta.chat.completions.parse(
        model=config.LLM_MODEL,
        messages=[
            {
                "role": "system",
                "content": """You are an expert storyteller and business case study writer for Reailize.
                            Your task is to create a list of unique, distinct success stories and product features 
                            that synthesizes insights from multiple framework approaches (STAR, HERO, PAS).
                            
                            You will be provided with:
                            1. Original extracted data from the source material
                            2. Individual framework stories (STAR, HERO, PAS)
                            3. Validation feedback highlighting strengths and areas for improvement
                            
                            Use the validation feedback to:
                            - Address any identified gaps or inconsistencies
                            - Emphasize the strongest elements from each framework
                            - Ensure each story is complete and compelling
                            
                            Create SEPARATE, UNIQUE stories for:
                            1. Each distinct customer implementation or project
                            2. Each unique product feature or capability demonstrated
                            3. Each different technical solution or innovation
                            4. Each separate business outcome or value proposition
                            
                            Each story should:
                            - Be self-contained and complete
                            - Focus on a specific customer, product, or technical achievement
                            - Include measurable outcomes and business value
                            - Be suitable for sales and marketing contexts
                            - Contain rich context and keywords for searchability
                            - Be compelling and professional for executive audiences
                            
                            IMPORTANT: Create multiple distinct stories rather than one unified narrative. 
                            Each story should represent a unique success case or product feature that can stand alone."""
            },
            {
                "role": "user",
                "content": f"""Create a list of unique success stories and product features using the validation insights provided:

ORIGINAL EXTRACTED DATA:
{context_text}

FRAMEWORK STORIES:
{framework_stories_text}

VALIDATION FEEDBACK:
{validation_feedback}

Please create separate, distinct success stories that:
- Each focuses on a unique customer project, product feature, or technical achievement
- Incorporates the validation feedback to address gaps and strengthen weak areas
- Combines the best validated insights from all framework stories
- Each tells a complete story from challenge to successful outcome
- Emphasizes measurable business value and technical achievements
- Includes relevant keywords and context for searchability
- Each is suitable for different use cases (sales presentations, case studies, reference materials)

OUTPUT REQUIREMENTS:
- Generate multiple unique stories (not one combined story)
- Each story should be distinct and self-contained
- Focus on different aspects: customer implementations, product features, technical innovations, business outcomes
- Ensure each story addresses any inconsistencies or gaps identified in validation
- Make each story compelling for executive audiences

Examples of separate stories you might create:
- Customer A's digital transformation project
- Product Feature X that solved industry challenges YZ
- Technical innovation Z that delivered measurable results
- Partnership collaboration that created new capabilities"""
            }
        ],
        response_format=SuccessStories
    )
    
    success_stories_response = response.choices[0].message.parsed
    success_stories_list = success_stories_response.SuccessStories if success_stories_response else []
    
    logger.info(f"Generated {len(success_stories_list)} unique success stories")
    
    # Create a formatted output for logging
    formatted_stories = ""
    for i, story in enumerate(success_stories_list, 1):
        formatted_stories += f"\n--- Story {i} ---\n{story}\n"
    
    # Log the validation and combination process
    process_summary = f"Validation and combination completed for {len(available_frameworks)} framework stories. Generated {len(success_stories_list)} unique success stories."
    
    return {
        "status": f"{process_summary}", 
        "combined_story": formatted_stories,
        "success_stories_list": success_stories_list
    }




async def ingest_to_weaviate_node(state: IngestionState):
    """Ingest the generated stories and raw data into Weaviate."""
    logger.info("Starting ingestion to Weaviate")
    
    try:
        # Create the ingest manager
        ingest_loader = create_ingest_manager()
        
        if state.star_story or state.hero_story or state.pas_story or state.success_stories_list:
            stories_dict = {
                "star_story": state.star_story or "",
                "hero_story": state.hero_story or "",
                "pas_story": state.pas_story or "",
                "success_stories_list": state.success_stories_list or [],
                #"combined_story": state.combined_story or ""
            }
            

        ExtractedData_dict = {
            "ExtractedData_list": state.ExtractedData_list or [],
            "source_document_link": state.source_document_link or state.pdf_path or "",
            "source_ppt_path": state.ppt_path or "",
            "source_pdf_path": state.pdf_path or ""
        }
        logger.info("Ingesting framework stories to Weaviate")
        await ingest_loader.ingest_framework_stories(
                stories_dict or {},
                source_metadata=ExtractedData_dict or {}
            )
        
        
        logger.info("Successfully completed ingestion to Weaviate")
        success_stories_count = len(state.success_stories_list) if state.success_stories_list else 0
        return {
            "success_ingestion": True, 
            "status": f"Successfully ingested data to Weaviate. Framework stories: {bool(state.star_story or state.hero_story or state.pas_story)}, Unique success stories: {success_stories_count}, Raw stories: {len(state.ExtractedData_list) if state.ExtractedData_list else 0}",
            "combined_story": state.combined_story,
            "success_stories_list": state.success_stories_list,
        }
        
    except Exception as e:
        logger.error(f"Error during Weaviate ingestion: {e}")
        return {
            "success_ingestion": False, 
            "status": f"Error during {state.pdf_path} ingestion to Weaviate: {str(e)}",
        }



ingestion_graph_builder = StateGraph(IngestionState,input_schema=DocIngestionRequest,output_schema=DocIngestionOutput)
    
ingestion_graph_builder.add_node("ProcessUploadedFiles", process_uploaded_files)
ingestion_graph_builder.add_node("ConvertToPDF", convert_to_pdf_node)
ingestion_graph_builder.add_node("ConvertPDFToImages", convert_pdf_to_images_node)
#ingestion_graph_builder.add_node("ExtractTextAndNotes", extract_text_and_notes_from_ppt_node)
ingestion_graph_builder.add_node("ExtractDataFromImage", extract_data_node)
ingestion_graph_builder.add_node("GenerateSTARStory", star_framework_node)
ingestion_graph_builder.add_node("GenerateHEROStory", hero_framework_node)
ingestion_graph_builder.add_node("GeneratePASStory", pas_framework_node)
ingestion_graph_builder.add_node("GenerateCombinedStory", validate_combined_story_node)
ingestion_graph_builder.add_node("IngestToWeaviate", ingest_to_weaviate_node)


ingestion_graph_builder.add_edge(START, "ProcessUploadedFiles")
ingestion_graph_builder.add_edge("ProcessUploadedFiles", "ConvertToPDF")
ingestion_graph_builder.add_edge("ConvertToPDF", "ConvertPDFToImages")
#ingestion_graph_builder.add_edge("ConvertPDFToImages", "ExtractTextAndNotes")
ingestion_graph_builder.add_edge("ConvertPDFToImages", "ExtractDataFromImage")
ingestion_graph_builder.add_edge("ExtractDataFromImage", "GenerateSTARStory")
ingestion_graph_builder.add_edge("GenerateSTARStory", "GenerateHEROStory")
ingestion_graph_builder.add_edge("GenerateHEROStory", "GeneratePASStory")
ingestion_graph_builder.add_edge("GeneratePASStory", "GenerateCombinedStory")
ingestion_graph_builder.add_edge("GenerateCombinedStory", "IngestToWeaviate")
ingestion_graph_builder.add_edge("IngestToWeaviate", END)

    
memory = MemorySaver()
ingestion_graph = ingestion_graph_builder.compile(checkpointer=memory)

# Cleanup function for proper thread pool shutdown
import atexit

def _cleanup_thread_pool():
    """Cleanup the global thread pool executor on module exit."""
    try:
        _thread_pool_executor.shutdown(wait=True)
        logger.debug("Thread pool executor shut down cleanly")
    except Exception as e:
        logger.warning(f"Error shutting down thread pool executor: {e}")

# Register cleanup function
atexit.register(_cleanup_thread_pool)


