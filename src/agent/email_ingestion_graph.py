import asyncio
import base64
import json
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import List, Optional, Tuple

import fitz
import httpx
from openai import AsyncOpenAI
from pptx import Presentation
from langgraph.graph import END, START, StateGraph

from src.agent.config import Config
from src.agent.graph_schemas import EmailClassificationResult, EmailIngestionOutput, EmailIngestionRequest, EmailIngestionState
from src.agent.logger import get_logger
from src.agent.excel_tracker import EmailClassificationExcelTracker

# Timeout constants for API calls
HTTP_TIMEOUT = httpx.Timeout(30.0)
LLM_TIMEOUT = 60.0  # seconds
TOKEN_REQUEST_TIMEOUT = 30.0  # seconds

logger = get_logger(__name__)
config = Config()

client = AsyncOpenAI(
    base_url=config.OPENAI_BASE_URL,
    api_key=config.OPENROUTER_API_KEY,
)

PROMPT_DIR = Path(__file__).parent / "prompts"
SYSTEM_PROMPT_FILE = PROMPT_DIR / "email_classification_system_prompt.txt"
USER_PROMPT_FILE = PROMPT_DIR / "email_classification_user_prompt.txt"


def _read_prompt(path: Path, fallback: str) -> str:
    try:
        if path.exists():
            return path.read_text(encoding="utf-8")
    except Exception as exc:
        logger.warning(f"Failed to read prompt file {path}: {exc}")
    return fallback


def _get_system_prompt() -> str:
    default_prompt = (
        "You classify inbound emails into business categories and return structured JSON only. "
        "Use the provided schema exactly. Confidence must be between 0.0 and 1.0."
    )
    return _read_prompt(SYSTEM_PROMPT_FILE, default_prompt)


def _get_user_prompt_template() -> str:
    default_template = (
        "Classify this inbound email inquiry.\n\n"
        "Subject:\n{email_subject}\n\n"
        "Body:\n{email_body}\n\n"
        "Attachment text:\n{attachment_text}\n\n"
        "Retrieved context:\n{retrieved_context}\n"
    )
    return _read_prompt(USER_PROMPT_FILE, default_template)


def _parse_ingested_email_content(email_content: str) -> Tuple[str, str, str]:
    subject = ""
    body = email_content
    attachment_text = ""

    if email_content.startswith("Subject:"):
        subject_end = email_content.find("\n\n")
        if subject_end != -1:
            subject_line = email_content[:subject_end]
            subject = subject_line.replace("Subject:", "", 1).strip()
            body = email_content[subject_end + 2 :]

    body_marker = "Body:\n"
    body_start = body.find(body_marker)
    if body_start != -1:
        body = body[body_start + len(body_marker) :]

    attachment_marker = "\n\n---\nExtracted Attachment Content\n---\n"
    attachment_idx = body.find(attachment_marker)
    if attachment_idx != -1:
        attachment_text = body[attachment_idx + len(attachment_marker) :].strip()
        body = body[:attachment_idx].strip()
    else:
        body = body.strip()

    return subject, body, attachment_text


async def get_access_token() -> Optional[str]:
    try:
        missing_fields = [
            key
            for key, value in {
                "CLIENT_ID": config.CLIENT_ID,
                "CLIENT_SECRET": config.CLIENT_SECRET,
                "TENANT_ID": config.TENANT_ID,
            }.items()
            if not value
        ]
        if missing_fields:
            logger.error(
                "Auth config missing required fields: %s. Check .env loading and runtime working directory. "
                "Make sure your .env file contains CLIENT_ID, CLIENT_SECRET, and TENANT_ID.",
                ", ".join(missing_fields),
            )
            return None

        token_url = f"https://login.microsoftonline.com/{config.TENANT_ID}/oauth2/v2.0/token"
        scope_value = " ".join(config.MSAL_SCOPE) if isinstance(config.MSAL_SCOPE, list) else str(config.MSAL_SCOPE)

        payload = {
            "client_id": config.CLIENT_ID,
            "client_secret": config.CLIENT_SECRET,
            "scope": scope_value,
            "grant_type": "client_credentials",
        }

        async def _request_token_via_httpx() -> dict:
            try:
                async with httpx.AsyncClient(timeout=HTTP_TIMEOUT) as async_client:
                    response = await asyncio.wait_for(
                        async_client.post(token_url, data=payload),
                        timeout=TOKEN_REQUEST_TIMEOUT
                    )

                if response.status_code != 200:
                    logger.error(
                        "Failed to get access token. Status: %s, Body: %s",
                        response.status_code,
                        response.text[:800],
                    )
                    return {}

                return response.json()
            except asyncio.TimeoutError:
                logger.error("Token request timed out after %s seconds", TOKEN_REQUEST_TIMEOUT)
                return {}
            except Exception as e:
                logger.error(f"HTTPX token request error: {e}")
                return {}

        async def _request_token_via_curl() -> dict:
            proc = await asyncio.create_subprocess_exec(
                "curl",
                "-sS",
                "-X",
                "POST",
                token_url,
                "-H",
                "Content-Type: application/x-www-form-urlencoded",
                "--data-urlencode",
                f"client_id={config.CLIENT_ID}",
                "--data-urlencode",
                f"client_secret={config.CLIENT_SECRET}",
                "--data-urlencode",
                f"scope={scope_value}",
                "--data-urlencode",
                "grant_type=client_credentials",
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            stdout, stderr = await proc.communicate()

            if proc.returncode != 0:
                logger.error("Curl token request failed: %s", stderr.decode("utf-8", errors="ignore")[:800])
                return {}

            try:
                return json.loads(stdout.decode("utf-8", errors="ignore") or "{}")
            except Exception as parse_exc:
                logger.error("Failed to parse curl token response: %s", parse_exc)
                return {}

        try:
            result = await _request_token_via_httpx()
        except Exception as exc:
            logger.warning("HTTPX token request failed, falling back to curl: %s", exc)
            result = await _request_token_via_curl()

        token = result.get("access_token")
        if token:
            return str(token)

        logger.error(
            "Failed to get access token. Error: %s, Description: %s",
            result.get("error", "unknown") if result else "no_result",
            result.get("error_description", "Unknown error") if result else "No result returned",
        )
        return None
    except Exception as exc:
        logger.error(f"Exception occurred while acquiring access token: {exc}")
        return None


async def get_email_messages(state: EmailIngestionState):
    logger.info(f"Get Access Token for email ID {state.email_id}")
    access_token = await get_access_token()
    if not access_token:
        return {"status": "failed: no access token", "email_content": ""}

    if not config.MAIL_USER:
        logger.error("MAIL_USER not configured in environment variables")
        return {"status": "failed: MAIL_USER not configured", "email_content": ""}

    url = (
        f"https://graph.microsoft.com/v1.0/users/{config.MAIL_USER}/messages/{state.email_id}"
        "?$select=id,conversationId,hasAttachments,subject,isRead,body,webLink,sender,from"
    )
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json",
        "Prefer": 'outlook.body-content-type="text"',
    }

    try:
        async with httpx.AsyncClient(timeout=HTTP_TIMEOUT) as async_client:
            response = await asyncio.wait_for(
                async_client.get(url, headers=headers),
                timeout=30.0
            )

        if response.status_code != 200:
            logger.error(f"Failed to fetch email message: {response.status_code} {response.text}")
            return {"status": "failed: could not fetch email", "email_content": ""}

        data = response.json()
        subject = data.get("subject", "")
        body_obj = data.get("body", {})
        body = body_obj.get("content", "") if body_obj else ""
        has_attachments = data.get("hasAttachments", False)
        sender_info = data.get("sender")
        sender_email = sender_info.get("emailAddress", {}).get("address") if sender_info else None

        email_content = f"Subject: {subject or 'No Subject'}\n\nBody:\n{body or '[No body content]'}"

        return {
            "status": "email retrieved",
            "email_content": email_content,
            "hasAttachments": has_attachments,
            "sender": sender_email,
        }
    except asyncio.TimeoutError:
        logger.error(f"Timeout fetching email message {state.email_id} - request exceeded 30 seconds")
        return {"status": "failed: email fetch timeout", "email_content": ""}
    except json.JSONDecodeError as exc:
        logger.error(f"Failed to parse email response JSON: {exc}")
        return {"status": "failed: invalid email response format", "email_content": ""}
    except Exception as exc:
        logger.error(f"Unexpected error while fetching email message: {exc}", exc_info=True)
        return {"status": f"failed: exception while fetching email - {type(exc).__name__}", "email_content": ""}


async def download_attachments(state: EmailIngestionState):
    if not state.hasAttachments:
        return {"status": "no attachments", "attachment_files": []}

    access_token = await get_access_token()
    if not access_token:
        return {"status": "failed: no access token for attachments", "attachment_files": []}

    if not config.MAIL_USER:
        logger.error("MAIL_USER not configured for attachment download")
        return {"status": "failed: MAIL_USER not configured", "attachment_files": []}

    url = f"https://graph.microsoft.com/v1.0/users/{config.MAIL_USER}/messages/{state.email_id}/attachments"
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json",
    }

    try:
        async with httpx.AsyncClient(timeout=HTTP_TIMEOUT) as async_client:
            response = await asyncio.wait_for(
                async_client.get(url, headers=headers),
                timeout=45.0
            )

        if response.status_code != 200:
            logger.error(f"Failed to fetch attachments: {response.status_code} {response.text}")
            return {"status": "failed: attachment fetch failed", "attachment_files": []}

        data = response.json()
    except asyncio.TimeoutError:
        logger.error(f"Timeout downloading attachments for email {state.email_id} - request exceeded 45 seconds")
        return {"status": "failed: attachment download timeout", "attachment_files": []}
    except json.JSONDecodeError as exc:
        logger.error(f"Failed to parse attachments response JSON: {exc}")
        return {"status": "failed: invalid attachments response format", "attachment_files": []}
    except Exception as exc:
        logger.error(f"Unexpected error while fetching attachments: {exc}", exc_info=True)
        return {"status": f"failed: exception while fetching attachments - {type(exc).__name__}", "attachment_files": []}

    files: List[Tuple[str, str]] = []
    for attachment in data.get("value", []):
        if attachment.get("@odata.type") != "#microsoft.graph.fileAttachment":
            continue

        name = attachment.get("name", "attachment.bin")
        content_b64 = attachment.get("contentBytes", "")
        if not content_b64:
            continue

        files.append((name, content_b64))

    return {
        "status": f"downloaded {len(files)} file attachments",
        "attachment_files": files,
    }


def _extract_text_from_pdf(raw_bytes: bytes, max_pages: int = 10) -> str:
    doc = fitz.open(stream=raw_bytes, filetype="pdf")
    try:
        chunks: List[str] = []
        page_count = min(len(doc), max_pages)
        for page_index in range(page_count):
            page_text = doc[page_index].get_text("text")
            chunks.append(str(page_text))
        return "\n".join(chunks).strip()
    finally:
        doc.close()


def _extract_text_from_pptx(raw_bytes: bytes, max_slides: int = 12) -> str:
    presentation = Presentation(BytesIO(raw_bytes))
    slide_chunks: List[str] = []

    for idx, slide in enumerate(list(presentation.slides)[:max_slides], start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            shape_text = getattr(shape, "text", None)
            if shape_text:
                parts.append(str(shape_text))

        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    parts.append(f"Notes: {notes}")
        except Exception:
            pass

        if parts:
            slide_chunks.append(f"Slide {idx}: {' '.join(parts)}")

    return "\n".join(slide_chunks).strip()


def _extract_text_by_filename(filename: str, raw_bytes: bytes) -> str:
    lower = filename.lower()

    if lower.endswith(".pdf"):
        return _extract_text_from_pdf(raw_bytes)

    if lower.endswith(".pptx"):
        return _extract_text_from_pptx(raw_bytes)

    if lower.endswith(".txt") or lower.endswith(".csv"):
        return raw_bytes.decode("utf-8", errors="ignore").strip()

    return ""


async def extract_attachment_text(state: EmailIngestionState):
    if not state.attachment_files:
        return {"status": "no attachment text to extract"}

    extracted_blocks: List[str] = []
    extracted_count = 0
    parse_failed_count = 0

    for filename, content_b64 in state.attachment_files:
        try:
            raw = base64.b64decode(content_b64)
            # Run blocking file parsing operations in a thread pool
            text = await asyncio.to_thread(_extract_text_by_filename, filename, raw)
            if text:
                extracted_blocks.append(f"Attachment: {filename}\n{text[:4000]}")
                extracted_count += 1
            else:
                extracted_blocks.append(f"Attachment: {filename}\n[No supported text extraction for this file type]")
        except Exception as exc:
            logger.warning(f"Attachment parse failed for {filename}: {exc}")
            parse_failed_count += 1
            extracted_blocks.append(f"Attachment: {filename}\n[Extraction failed]")

    attachment_context = "\n\n".join(extracted_blocks).strip()

    combined_email_content = state.email_content or ""
    if attachment_context:
        combined_email_content += "\n\n---\nExtracted Attachment Content\n---\n"
        combined_email_content += attachment_context

    return {
        "email_content": combined_email_content,
        "status": (
            f"attachment extraction complete: {extracted_count} parsed, "
            f"{parse_failed_count} failed, {len(state.attachment_files)} total"
        ),
    }


async def classify_email(state: EmailIngestionState):
    if not state.email_content:
        fallback = EmailClassificationResult(
            date_of_contact="",
            action="disqualify",
            company_name="",
            company_type="unknown",
            operation_countries=[],
            company_presence=[],
            current_projects=[],
            source="",
            email="",
            contact_name="",
            contact_last_name="",
            salesperson="none",
            confidence=0.0,
        )
        fallback_payload = fallback.model_dump()
        fallback_payload["error"] = "classification_skipped: missing email_content"
        return {
            "classification": fallback_payload,
            "status": "classification skipped: no email content",
        }

    if not config.OPENROUTER_API_KEY:
        logger.error("OPENROUTER_API_KEY not configured - cannot classify email")
        fallback = EmailClassificationResult(
            date_of_contact="",
            action="disqualify",
            company_name="",
            company_type="unknown",
            operation_countries=[],
            company_presence=[],
            current_projects=[],
            source="",
            email="",
            contact_name="",
            contact_last_name="",
            salesperson="none",
            confidence=0.0,
        )
        fallback_payload = fallback.model_dump()
        fallback_payload["error"] = "classification_skipped: missing API key"
        return {
            "classification": fallback_payload,
            "status": "classification skipped: API key not configured",
        }

    email_subject, email_body, attachment_text = _parse_ingested_email_content(state.email_content)
    system_prompt = _get_system_prompt()
    user_prompt_template = _get_user_prompt_template()

    user_prompt = user_prompt_template.format(
        email_subject=email_subject or "",
        email_body=email_body or "",
        attachment_text=attachment_text or "",
        retrieved_context="",
    )

    try:
        response = await asyncio.wait_for(
            client.beta.chat.completions.parse(
                model=config.LLM_MODEL,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt},
                ],
                response_format=EmailClassificationResult,
            ),
            timeout=LLM_TIMEOUT
        )

        parsed = response.choices[0].message.parsed
        if not parsed:
            raise ValueError("No parsed response returned by model")

        result = parsed.model_dump()
        result["confidence"] = max(0.0, min(1.0, float(result.get("confidence", 0.0))))

        return {
            "classification": result,
            "status": "email classified successfully",
        }
    except asyncio.TimeoutError:
        logger.error(f"Email classification timeout after {LLM_TIMEOUT}s for email {state.email_id}")
        fallback = EmailClassificationResult(
            date_of_contact="",
            action="disqualify",
            company_name="",
            company_type="unknown",
            operation_countries=[],
            company_presence=[],
            current_projects=[],
            source="",
            email="",
            contact_name="",
            contact_last_name="",
            salesperson="none",
            confidence=0.0,
        )
        fallback_payload = fallback.model_dump()
        fallback_payload["error"] = "classification_failed: timeout"
        return {
            "classification": fallback_payload,
            "status": f"email classification timed out after {LLM_TIMEOUT}s; fallback returned",
        }
    except Exception as exc:
        logger.error(f"Email classification failed in ingestion graph: {type(exc).__name__}: {exc}", exc_info=True)
        fallback = EmailClassificationResult(
            date_of_contact="",
            action="disqualify",
            company_name="",
            company_type="unknown",
            operation_countries=[],
            company_presence=[],
            current_projects=[],
            source="",
            email="",
            contact_name="",
            contact_last_name="",
            salesperson="none",
            confidence=0.0,
        )
        fallback_payload = fallback.model_dump()
        fallback_payload["error"] = f"classification_failed: {type(exc).__name__}"
        return {
            "classification": fallback_payload,
            "status": f"email classification failed ({type(exc).__name__}); fallback returned",
        }


def _append_to_excel_sync(
    thread_id: str,
    created_at: str,
    email_id: str,
    sender: str,
    email_content: str,
    classification: dict,
    status: str,
) -> tuple[bool, str]:
    """Synchronous function to handle all Excel operations in a thread."""
    try:
        tracker = EmailClassificationExcelTracker()
        success = tracker.append_email(
            thread_id=thread_id,
            created_at=created_at,
            email_id=email_id,
            sender=sender,
            email_content=email_content,
            classification=classification,
            status=status,
        )
        return success, str(tracker.excel_file)
    except Exception as exc:
        logger.error(f"Excel append error: {exc}", exc_info=True)
        return False, str(exc)


async def append_to_excel(state: EmailIngestionState):
    """Append classified email to Excel file (optional step that doesn't fail the workflow)."""
    try:
        # Get thread_id from conversation_id or generate one
        thread_id = state.conversation_id or "unknown"
        
        # Get current timestamp if not available
        created_at = datetime.now().isoformat()
        
        # Run ALL synchronous Excel operations (including initialization) in a thread pool
        # This prevents blocking the event loop with file I/O and directory creation
        success, excel_path = await asyncio.wait_for(
            asyncio.to_thread(
                _append_to_excel_sync,
                thread_id=thread_id,
                created_at=created_at,
                email_id=state.email_id or "",
                sender=state.sender or "",
                email_content=state.email_content or "",
                classification=state.classification or {},
                status="processed",
            ),
            timeout=30.0  # 30 second timeout for Excel operations
        )
        
        if success:
            logger.info(f"Successfully appended email to Excel: {excel_path}")
            return {
                "status": f"email appended to Excel: {excel_path}",
            }
        else:
            logger.warning(f"Excel append returned False: {excel_path}")
            return {
                "status": "excel append completed with warnings",
            }
    except asyncio.TimeoutError:
        logger.error("Excel append operation timed out after 30s")
        return {
            "status": "excel append timed out (non-fatal)",
        }
    except Exception as exc:
        logger.error(f"Failed to append to Excel: {exc}", exc_info=True)
        return {
            "status": f"excel append failed: {type(exc).__name__} (non-fatal)",
        }


email_ingestion_graph_builder = StateGraph(
    EmailIngestionState,
    input_schema=EmailIngestionRequest,
    output_schema=EmailIngestionOutput,
)

email_ingestion_graph_builder.add_node("get_email_messages", get_email_messages)
email_ingestion_graph_builder.add_node("download_attachments", download_attachments)
email_ingestion_graph_builder.add_node("extract_attachment_text", extract_attachment_text)
email_ingestion_graph_builder.add_node("classify_email", classify_email)
email_ingestion_graph_builder.add_node("append_to_excel", append_to_excel)

email_ingestion_graph_builder.add_edge(START, "get_email_messages")
email_ingestion_graph_builder.add_edge("get_email_messages", "download_attachments")
email_ingestion_graph_builder.add_edge("download_attachments", "extract_attachment_text")
email_ingestion_graph_builder.add_edge("extract_attachment_text", "classify_email")
email_ingestion_graph_builder.add_edge("classify_email", "append_to_excel")
email_ingestion_graph_builder.add_edge("append_to_excel", END)

email_ingestion_graph = email_ingestion_graph_builder.compile()
